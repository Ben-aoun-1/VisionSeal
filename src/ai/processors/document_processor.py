"""
Enhanced Document Processing for VisionSeal RAG System
Supports PDF, DOCX, and PPTX with robust error handling
"""
import os
import io
import json
import re
from typing import List, Tuple, Dict, Any, Optional
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor
import tempfile

# Document processing libraries
from docx import Document
import fitz  # PyMuPDF
from pptx import Presentation

# OCR libraries (optional)
try:
    import pytesseract
    from PIL import Image
    import numpy as np
    import cv2
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

# Text splitting
try:
    from langchain.text_splitter import RecursiveCharacterTextSplitter
    LANGCHAIN_AVAILABLE = True
except ImportError:
    LANGCHAIN_AVAILABLE = False

from core.logging.setup import get_logger
from core.exceptions.handlers import AIProcessingException

logger = get_logger("document_processor")


class DocumentProcessor:
    """Enhanced document processor with multi-format support"""
    
    def __init__(
        self,
        chunk_size: int = 1500,
        chunk_overlap: int = 200,
        enable_ocr: bool = False
    ):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.enable_ocr = enable_ocr and OCR_AVAILABLE
        
        if self.enable_ocr and not OCR_AVAILABLE:
            logger.warning("OCR requested but dependencies not available")
            self.enable_ocr = False

    async def process_document(
        self, 
        file_path: str, 
        file_type: Optional[str] = None
    ) -> List[Tuple[str, str, str]]:
        """
        Process document and return chunks
        Returns: List of (section, content, type) tuples
        """
        try:
            if not os.path.exists(file_path):
                raise FileNotFoundError(f"Document not found: {file_path}")
            
            # Determine file type
            if not file_type:
                file_type = Path(file_path).suffix.lower()
            
            logger.info(f"Processing document: {file_path} (type: {file_type})")
            
            # Extract sections based on file type
            if file_type in ['.pdf']:
                sections = await self._extract_pdf_sections(file_path)
            elif file_type in ['.docx', '.doc']:
                sections = await self._extract_docx_sections(file_path)
            elif file_type in ['.pptx', '.ppt']:
                sections = await self._extract_pptx_sections(file_path)
            else:
                raise AIProcessingException(f"Unsupported file type: {file_type}")
            
            # Chunk the sections
            chunks = await self._chunk_sections(sections)
            
            logger.info(f"Document processed: {len(chunks)} chunks created")
            return chunks
            
        except Exception as e:
            logger.error(f"Document processing failed: {str(e)}")
            raise AIProcessingException(f"Document processing failed: {str(e)}")

    async def _extract_pdf_sections(self, file_path: str) -> List[Tuple[str, List[Tuple[str, Any]]]]:
        """Extract sections from PDF with optional OCR"""
        try:
            def process_page(page_info):
                page, page_num = page_info
                try:
                    # Try text extraction first
                    text = page.get_text("text").strip()
                    
                    # If minimal text and OCR enabled, try OCR
                    if len(text) < 50 and self.enable_ocr:
                        text = self._ocr_page(page)
                    
                    return f"Page {page_num + 1}", text
                except Exception as e:
                    logger.warning(f"Error processing page {page_num + 1}: {str(e)}")
                    return f"Page {page_num + 1}", ""

            with fitz.open(file_path) as doc:
                # Process pages in parallel
                with ThreadPoolExecutor(max_workers=4) as executor:
                    page_results = list(executor.map(
                        process_page,
                        [(page, i) for i, page in enumerate(doc)]
                    ))
                
                sections = []
                for section_title, content in page_results:
                    if content.strip():
                        sections.append((section_title, [('text', content)]))
                
                return sections
                
        except Exception as e:
            logger.error(f"PDF extraction failed: {str(e)}")
            raise AIProcessingException(f"PDF extraction failed: {str(e)}")

    def _ocr_page(self, page) -> str:
        """Perform OCR on a PDF page"""
        if not self.enable_ocr:
            return ""
        
        try:
            # Get page as image
            pix = page.get_pixmap(dpi=300)
            img_data = pix.tobytes("png")
            img = Image.open(io.BytesIO(img_data))
            
            # Convert to numpy array for OpenCV processing
            img_np = np.array(img)
            gray = cv2.cvtColor(img_np, cv2.COLOR_RGB2GRAY)
            
            # Apply image preprocessing
            denoised = cv2.fastNlMeansDenoising(gray, h=30)
            enhanced = cv2.createCLAHE(clipLimit=3.0, tileGridSize=(8,8)).apply(denoised)
            
            # Perform OCR
            custom_config = r'--oem 3 --psm 6 -l fra+eng'
            text = pytesseract.image_to_string(
                Image.fromarray(enhanced),
                config=custom_config
            )
            
            # Clean up text
            return re.sub(r'(?<!\n)\n(?!\n)', ' ', text).strip()
            
        except Exception as e:
            logger.warning(f"OCR failed: {str(e)}")
            return ""

    async def _extract_docx_sections(self, file_path: str) -> List[Tuple[str, List[Tuple[str, Any]]]]:
        """Extract sections from DOCX with table handling"""
        try:
            doc = Document(file_path)
            sections = []
            current_section = "Document"
            current_content = []

            # Process paragraphs
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if not text:
                    continue
                    
                if self._is_heading(paragraph):
                    # Save previous section
                    if current_content:
                        sections.append((current_section, current_content.copy()))
                    current_section = text
                    current_content = []
                else:
                    current_content.append(('text', text))
            
            # Process tables
            for table in doc.tables:
                table_data = self._table_to_json(table)
                if table_data:
                    current_content.append(('table', table_data))
            
            # Add final section
            if current_content:
                sections.append((current_section, current_content))

            return sections
            
        except Exception as e:
            logger.error(f"DOCX extraction failed: {str(e)}")
            raise AIProcessingException(f"DOCX extraction failed: {str(e)}")

    def _is_heading(self, paragraph) -> bool:
        """Determine if paragraph is a heading"""
        return (
            any(run.bold for run in paragraph.runs) or
            any(run.font.size and run.font.size.pt >= 12 for run in paragraph.runs if run.font.size) or
            paragraph.style.name.startswith('Heading')
        )

    def _table_to_json(self, table) -> List[Dict[str, str]]:
        """Convert Word table to JSON structure"""
        try:
            table_data = []
            # Use first row as headers if available
            headers = [cell.text.strip() for cell in table.rows[0].cells] if table.rows else []
            
            for i, row in enumerate(table.rows):
                # Skip header row
                if i == 0 and headers:
                    continue
                    
                row_data = {}
                for j, cell in enumerate(row.cells):
                    key = headers[j] if j < len(headers) and headers[j] else f"col_{j}"
                    row_data[key] = cell.text.strip()
                
                if any(row_data.values()):  # Only add non-empty rows
                    table_data.append(row_data)
            
            return table_data
            
        except Exception as e:
            logger.warning(f"Table conversion failed: {str(e)}")
            return []

    async def _extract_pptx_sections(self, file_path: str) -> List[Tuple[str, List[Tuple[str, Any]]]]:
        """Extract sections from PPTX"""
        try:
            presentation = Presentation(file_path)
            sections = []

            for i, slide in enumerate(presentation.slides):
                slide_title = f"Slide {i + 1}"
                content = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content.append(('text', shape.text.strip()))
                
                if content:
                    sections.append((slide_title, content))

            return sections
            
        except Exception as e:
            logger.error(f"PPTX extraction failed: {str(e)}")
            raise AIProcessingException(f"PPTX extraction failed: {str(e)}")

    async def _chunk_sections(self, sections: List[Tuple[str, List[Tuple[str, Any]]]]) -> List[Tuple[str, str, str]]:
        """Chunk sections into manageable pieces"""
        try:
            chunks = []
            
            for section_title, content_items in sections:
                current_chunk = f"{section_title}\n\n" if section_title else ""
                
                for item_type, item_content in content_items:
                    if item_type == 'text':
                        current_chunk += item_content + "\n\n"
                    elif item_type == 'table':
                        # Convert table to text representation
                        table_text = self._table_to_text(item_content)
                        
                        # Handle large tables separately
                        if len(table_text) > self.chunk_size * 0.8:
                            # Save current chunk if not empty
                            if current_chunk.strip():
                                chunks.append((section_title, current_chunk.strip(), "text"))
                                current_chunk = ""
                            # Add table as separate chunk
                            chunks.append((section_title, table_text, "table"))
                        else:
                            current_chunk += table_text + "\n\n"
                
                # Split large chunks
                if current_chunk.strip():
                    if LANGCHAIN_AVAILABLE:
                        split_chunks = self._split_with_langchain(current_chunk)
                        for chunk in split_chunks:
                            chunks.append((section_title, chunk.strip(), "text"))
                    else:
                        # Simple splitting fallback
                        chunks.append((section_title, current_chunk.strip(), "text"))

            return chunks
            
        except Exception as e:
            logger.error(f"Chunking failed: {str(e)}")
            raise AIProcessingException(f"Chunking failed: {str(e)}")

    def _table_to_text(self, table_data: List[Dict[str, str]]) -> str:
        """Convert table data to text representation"""
        if not table_data:
            return ""
        
        try:
            # Create JSON representation with clear markers
            table_json = json.dumps(table_data, indent=2, ensure_ascii=False)
            return f"[TABLE_DATA]\n{table_json}\n[/TABLE_DATA]"
        except Exception:
            # Fallback to simple text representation
            text_lines = []
            for row in table_data:
                row_text = " | ".join(f"{k}: {v}" for k, v in row.items() if v)
                if row_text:
                    text_lines.append(row_text)
            return "\n".join(text_lines)

    def _split_with_langchain(self, text: str) -> List[str]:
        """Split text using LangChain text splitter"""
        try:
            splitter = RecursiveCharacterTextSplitter(
                chunk_size=self.chunk_size,
                chunk_overlap=self.chunk_overlap,
                separators=["\n\n", "\n", " ", ""]
            )
            return splitter.split_text(text)
        except Exception as e:
            logger.warning(f"LangChain splitting failed: {str(e)}")
            # Simple fallback splitting
            return [text[i:i+self.chunk_size] for i in range(0, len(text), self.chunk_size - self.chunk_overlap)]


# Factory function for easy instantiation
def create_document_processor(
    chunk_size: int = 1500,
    chunk_overlap: int = 200,
    enable_ocr: bool = False
) -> DocumentProcessor:
    """Create a document processor with specified configuration"""
    return DocumentProcessor(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        enable_ocr=enable_ocr
    )